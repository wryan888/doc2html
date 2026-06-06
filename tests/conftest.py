"""共用 fixtures：在 tmp_path 動態產生各格式測試檔。

二進位格式（docx/xlsx/pptx/pdf）若缺對應套件就 skip 相關測試，
讓核心測試在零依賴環境也能跑。
"""

import io
import json
import struct
import zlib

import pytest


def _make_png(width=2, height=2):
    """程式產生一張合法的 RGB PNG（紅色），含正確 CRC，給內嵌圖片測試用。"""

    def chunk(typ, data):
        body = typ + data
        crc = zlib.crc32(body) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + body + struct.pack(">I", crc)

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)  # 8-bit RGB
    raw = b"".join(b"\x00" + b"\xff\x00\x00" * width for _ in range(height))
    idat = zlib.compress(raw)
    return sig + chunk(b"IHDR", ihdr) + chunk(b"IDAT", idat) + chunk(b"IEND", b"")


PNG_1x1 = _make_png()


@pytest.fixture
def txt_file(tmp_path):
    p = tmp_path / "sample.txt"
    p.write_text("第一段。\n同段第二行。\n\n第二段。", encoding="utf-8")
    return p


@pytest.fixture
def csv_file(tmp_path):
    p = tmp_path / "sample.csv"
    p.write_text("名稱,價格\n蘋果,30\n香蕉,15", encoding="utf-8")
    return p


@pytest.fixture
def tsv_file(tmp_path):
    p = tmp_path / "sample.tsv"
    p.write_text("a\tb\n1\t2", encoding="utf-8")
    return p


@pytest.fixture
def json_file(tmp_path):
    p = tmp_path / "sample.json"
    p.write_text(
        json.dumps(
            {"k": "v", "items": [1, 2], "ok": True, "n": None},
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    return p


@pytest.fixture
def html_file(tmp_path):
    p = tmp_path / "sample.html"
    p.write_text(
        "<html><head><title>標題</title><style>p{color:red}</style></head>"
        "<body><article><h1>H</h1><p>內文<script>x()</script></p>"
        "</article></body></html>",
        encoding="utf-8",
    )
    return p


@pytest.fixture
def docx_file(tmp_path):
    docx = pytest.importorskip("docx")
    p = tmp_path / "sample.docx"
    d = docx.Document()
    d.add_heading("報告標題", level=1)
    d.add_paragraph("內文")
    d.add_paragraph("項目一", style="List Bullet")
    d.add_paragraph("項目二", style="List Bullet")
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "A"
    t.cell(0, 1).text = "B"
    t.cell(1, 0).text = "1"
    t.cell(1, 1).text = "2"
    d.save(p)
    return p


@pytest.fixture
def docx_rich_file(tmp_path):
    """含超連結與內嵌圖片的 DOCX。"""
    docx = pytest.importorskip("docx")
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    from docx.oxml.shared import OxmlElement, qn

    p = tmp_path / "rich.docx"
    d = docx.Document()
    para = d.add_paragraph("前往 ")

    # 手動加一個外部超連結 run
    r_id = para.part.relate_to(
        "https://example.com", RT.HYPERLINK, is_external=True
    )
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    run = OxmlElement("w:r")
    t = OxmlElement("w:t")
    t.text = "範例網站"
    run.append(t)
    hyperlink.append(run)
    para._p.append(hyperlink)

    # 內嵌圖片（自成一段）
    d.add_picture(io.BytesIO(PNG_1x1))
    d.save(p)
    return p


@pytest.fixture
def docx_anchor_file(tmp_path):
    """含內部錨點超連結（w:anchor，無 r:id）的 DOCX。"""
    docx = pytest.importorskip("docx")
    from docx.oxml.shared import OxmlElement, qn

    p = tmp_path / "anchor.docx"
    d = docx.Document()
    para = d.add_paragraph("跳到 ")
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("w:anchor"), "_Toc123")
    run = OxmlElement("w:r")
    t = OxmlElement("w:t")
    t.text = "第一章"
    run.append(t)
    hyperlink.append(run)
    para._p.append(hyperlink)
    d.save(p)
    return p


@pytest.fixture
def xlsx_file(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    p = tmp_path / "sample.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "庫存"
    ws.append(["名稱", "價格"])
    ws.append(["蘋果", 30])
    wb.save(p)
    return p


@pytest.fixture
def xlsx_no_header_file(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    p = tmp_path / "nohdr.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append([1, 2, 3])
    ws.append([4, 5, 6])
    wb.save(p)
    return p


@pytest.fixture
def pptx_file(tmp_path):
    pptx = pytest.importorskip("pptx")
    p = tmp_path / "sample.pptx"
    prs = pptx.Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "投影片標題"
    s.placeholders[1].text = "重點一\n重點二"
    prs.save(p)
    return p


@pytest.fixture
def pptx_image_file(tmp_path):
    """含內嵌圖片的 PPTX。"""
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    p = tmp_path / "img.pptx"
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "圖片頁"
    slide.shapes.add_picture(
        io.BytesIO(PNG_1x1), Inches(1), Inches(1), Inches(1), Inches(1)
    )
    prs.save(p)
    return p


@pytest.fixture
def pdf_table_file(tmp_path):
    """含一個有框線表格的 PDF（pdfplumber 可偵測）。"""
    pytest.importorskip("pdfplumber")
    pytest.importorskip("reportlab")
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, TableStyle
    from reportlab.platypus import Table as RLTable

    p = tmp_path / "table.pdf"
    doc = SimpleDocTemplate(str(p))
    data = [["Name", "Price"], ["Apple", "30"], ["Banana", "15"]]
    t = RLTable(data)
    t.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 1, colors.black)]))
    doc.build([t])
    return p


@pytest.fixture
def pdf_mixed_file(tmp_path):
    """含「標題 → 表格 → 內文」的 PDF，用來驗證交錯排序。"""
    pytest.importorskip("pdfplumber")
    pytest.importorskip("reportlab")
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, TableStyle
    from reportlab.platypus import Table as RLTable

    p = tmp_path / "mixed.pdf"
    styles = getSampleStyleSheet()
    t = RLTable([["Col", "Val"], ["x", "1"]])
    t.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 1, colors.black)]))
    story = [
        Paragraph("TopHeading", styles["Title"]),
        Spacer(1, 12),
        t,
        Spacer(1, 12),
        Paragraph("Bottom body text here.", styles["Normal"]),
    ]
    SimpleDocTemplate(str(p)).build(story)
    return p


@pytest.fixture
def pdf_file(tmp_path):
    pytest.importorskip("reportlab")
    pytest.importorskip("pdfplumber")
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    p = tmp_path / "sample.pdf"
    c = canvas.Canvas(str(p), pagesize=letter)
    c.setFont("Helvetica-Bold", 18)
    c.drawString(72, 720, "PDF Heading")
    c.setFont("Helvetica", 11)
    c.drawString(72, 690, "Body paragraph text.")
    c.save()
    return p
