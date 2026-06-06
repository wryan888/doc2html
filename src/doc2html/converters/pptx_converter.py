"""PowerPoint (.pptx) 轉換器：每張投影片轉成一個 <section>。

抽出標題、條列文字、表格與圖片的替代文字（alt text）。
"""

from __future__ import annotations

from typing import BinaryIO

from .._base_converter import DocumentConverter, DocumentConverterResult
from .._exceptions import MissingDependencyException
from .._html_builder import escape, table
from .._stream_info import StreamInfo

_PPTX_EXT = ".pptx"
_PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


class PptxConverter(DocumentConverter):
    priority = 5.0

    def accepts(self, file_stream: BinaryIO, stream_info: StreamInfo) -> bool:
        ext = (stream_info.extension or "").lower()
        mime = (stream_info.mimetype or "").lower()
        return ext == _PPTX_EXT or mime == _PPTX_MIME

    def convert(
        self, file_stream: BinaryIO, stream_info: StreamInfo
    ) -> DocumentConverterResult:
        try:
            from pptx import Presentation  # noqa: PLC0415
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415
        except ImportError as exc:
            raise MissingDependencyException(
                "PptxConverter 需要 python-pptx，請執行："
                "pip install 'doc2html[pptx]'"
            ) from exc

        prs = Presentation(file_stream)
        parts: list[str] = []
        first_title = None

        for index, slide in enumerate(prs.slides, start=1):
            section: list[str] = [f'<section class="slide" id="slide-{index}">']
            slide_title = _slide_title(slide)
            if slide_title:
                section.append(f"<h2>{escape(slide_title)}</h2>")
                if first_title is None:
                    first_title = slide_title
            else:
                section.append(f"<h2>投影片 {index}</h2>")

            for shape in slide.shapes:
                if shape == _title_shape(slide):
                    continue  # 標題已輸出
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    alt = _picture_alt(shape) or "圖片"
                    section.append(f'<p><em>[圖片：{escape(alt)}]</em></p>')
                elif shape.has_table:
                    section.append(_table_html(shape.table))
                elif shape.has_text_frame:
                    section.append(_text_frame_html(shape.text_frame))

            # 講者備忘稿
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    section.append(
                        f'<aside class="notes"><strong>備忘稿：</strong>'
                        f"{escape(notes)}</aside>"
                    )
            section.append("</section>")
            parts.append("\n".join(section))

        return DocumentConverterResult(
            "\n".join(parts), title=first_title or stream_info.filename
        )


def _title_shape(slide):
    try:
        return slide.shapes.title
    except (AttributeError, KeyError):
        return None


def _slide_title(slide):
    shape = _title_shape(slide)
    if shape is not None and shape.has_text_frame:
        text = shape.text_frame.text.strip()
        return text or None
    return None


def _picture_alt(shape):
    try:
        return shape._element._nvXxPr.cNvPr.get("descr") or shape.name
    except AttributeError:
        return shape.name


def _text_frame_html(text_frame) -> str:
    items = []
    for para in text_frame.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if text:
            items.append(f"<li>{escape(text)}</li>")
    if not items:
        return ""
    return "<ul>" + "".join(items) + "</ul>"


def _table_html(tbl) -> str:
    rows = [[cell.text for cell in row.cells] for row in tbl.rows]
    if not rows:
        return ""
    header, *body = rows
    return table(body, header=header)
