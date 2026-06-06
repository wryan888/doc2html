"""產生各格式的測試樣本檔，用來驗證 doc2html。"""
import json
import os

HERE = os.path.dirname(__file__)


def p(name):
    return os.path.join(HERE, name)


# --- 純文字 / CSV / JSON / HTML (內建，不需第三方) ---
with open(p("sample.txt"), "w", encoding="utf-8") as f:
    f.write("第一段。\n這是同一段的第二行。\n\n第二段獨立成段。")

with open(p("sample.csv"), "w", encoding="utf-8") as f:
    f.write("名稱,價格,庫存\n蘋果,30,100\n香蕉,15,250\n芒果,55,40")

with open(p("sample.json"), "w", encoding="utf-8") as f:
    json.dump(
        {"標題": "庫存", "項目": [{"名稱": "蘋果", "價格": 30},
         {"名稱": "香蕉", "價格": 15}], "啟用": True, "備註": None},
        f, ensure_ascii=False, indent=2,
    )

with open(p("sample.html"), "w", encoding="utf-8") as f:
    f.write("<!DOCTYPE html><html><head><title>原始標題</title>"
            "<style>p{color:red}</style></head><body>"
            "<article><h1>文章標題</h1><p>段落內容 <strong>粗體</strong>。</p>"
            "<script>alert('x')</script></article></body></html>")

# --- DOCX ---
try:
    import docx
    d = docx.Document()
    d.add_heading("報告標題", level=1)
    d.add_paragraph("這是一段內文。")
    para = d.add_paragraph()
    para.add_run("粗體").bold = True
    para.add_run(" 與 ")
    para.add_run("斜體").italic = True
    d.add_heading("子章節", level=2)
    d.add_paragraph("項目一", style="List Bullet")
    d.add_paragraph("項目二", style="List Bullet")
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "欄A"
    t.cell(0, 1).text = "欄B"
    t.cell(1, 0).text = "1"
    t.cell(1, 1).text = "2"
    d.save(p("sample.docx"))
    print("docx OK")
except Exception as e:
    print("docx 跳過:", e)

# --- XLSX ---
try:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "庫存"
    ws.append(["名稱", "價格", "庫存"])
    ws.append(["蘋果", 30, 100])
    ws.append(["香蕉", 15, 250])
    ws2 = wb.create_sheet("摘要")
    ws2.append(["總項目", 2])
    wb.save(p("sample.xlsx"))
    print("xlsx OK")
except Exception as e:
    print("xlsx 跳過:", e)

# --- PPTX ---
try:
    from pptx import Presentation
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "簡報標題"
    s1.placeholders[1].text = "重點一\n重點二"
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "第二張"
    prs.save(p("sample.pptx"))
    print("pptx OK")
except Exception as e:
    print("pptx 跳過:", e)

# --- PDF (用 reportlab 若有，否則跳過) ---
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(p("sample.pdf"), pagesize=letter)
    c.setFont("Helvetica-Bold", 18)
    c.drawString(72, 720, "PDF Heading")
    c.setFont("Helvetica", 11)
    c.drawString(72, 690, "This is a body paragraph in the PDF sample.")
    c.save()
    print("pdf OK")
except Exception as e:
    print("pdf 跳過 (需 reportlab):", e)

print("樣本產生完成")
